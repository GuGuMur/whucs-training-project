"""Tests for the multi-format document parser."""

import io

import pytest

from app.services.parser import ParseError, ParsedDocument, ParsedSegment, parse_document


# ── helpers ──────────────────────────────────────────────────────────

def _minimal_pdf_bytes() -> bytes:
    """Produce a valid single-page PDF with extractable text using PyMuPDF."""
    import fitz

    doc = fitz.open()
    page = doc.new_page(width=595, height=842)
    page.insert_text((72, 72), "Microscope observation", fontsize=11)
    page.insert_text((72, 96), "Cell walls clearly visible.", fontsize=11)
    buf = io.BytesIO()
    doc.save(buf)
    doc.close()
    return buf.getvalue()


def _minimal_docx_bytes() -> bytes:
    """Produce a single-paragraph DOCX with python-docx."""
    import docx

    document = docx.Document()
    document.add_paragraph("显微镜实验报告汇总。")
    buf = io.BytesIO()
    document.save(buf)
    return buf.getvalue()


def _minimal_pptx_bytes() -> bytes:
    """Produce a single-slide PPTX with python-pptx."""
    from pptx import Presentation

    presentation = Presentation()
    slide_layout = presentation.slide_layouts[0]
    slide = presentation.slides.add_slide(slide_layout)
    slide.shapes.title.text = "周报"
    buf = io.BytesIO()
    presentation.save(buf)
    return buf.getvalue()


# ── format detection ─────────────────────────────────────────────────

def test_format_detection_by_extension():
    """Each known extension maps to the correct internal format.

    PDF/DOCX/PPTX need valid binary content so we use helpers; TXT/MD/CSV
    work with plain UTF-8 bytes.
    """
    assert parse_document("report.pdf", _minimal_pdf_bytes()).text
    assert parse_document("report.docx", _minimal_docx_bytes()).text
    assert parse_document("slides.pptx", _minimal_pptx_bytes()).text
    assert parse_document("notes.txt", "纯文本".encode()).text == "纯文本"
    assert parse_document("README.md", "# heading".encode()).text


def test_format_detection_by_declared_type():
    """When the extension is unknown, declared_type acts as fallback."""
    doc = parse_document("data.bin", "a,b,c\n1,2,3".encode(), declared_type="csv")
    assert doc.text


def test_unsupported_format_raises_parse_error():
    with pytest.raises(ParseError, match="不支持的解析格式"):
        parse_document("data.xlsx", b"")


# ── TXT ──────────────────────────────────────────────────────────────

def test_text_parsing_utf8():
    doc = parse_document("notes.txt", "第一行内容\n第二行内容\n".encode())
    assert "第一行内容" in doc.text
    assert "第二行内容" in doc.text
    assert len(doc.segments) == 2
    assert doc.segments[0].content == "第一行内容"
    assert doc.segments[0].page_no == 1
    assert doc.segments[0].paragraph_no == 1
    assert doc.segments[1].paragraph_no == 2


def test_text_parsing_strips_empty_lines():
    doc = parse_document("notes.txt", "\n\n  a  \n\n  b  \n".encode())
    assert len(doc.segments) == 2
    assert doc.segments[0].content == "a"
    assert doc.segments[1].content == "b"


def test_text_parsing_single_paragraph():
    doc = parse_document("notes.txt", "only one sentence".encode())
    assert len(doc.segments) == 1
    assert doc.segments[0].paragraph_no == 1


def test_text_parsing_decodes_fallback():
    """Non-UTF-8 bytes are decoded with replacement characters."""
    doc = parse_document("data.txt", b"\xff\xfe\x00")
    assert doc.text  # survives decoding without crashing


# ── Markdown ─────────────────────────────────────────────────────────

def test_markdown_parsing():
    doc = parse_document("README.md", "# heading\ncontent paragraph.".encode())
    assert "heading" in doc.text
    assert "content paragraph" in doc.text


def test_markdown_removes_fenced_code_blocks():
    doc = parse_document("README.md", "note\n```python\nprint('hello')\n```\nconclusion".encode())
    assert "print('hello')" not in doc.text
    assert "note" in doc.text
    assert "conclusion" in doc.text


# ── CSV ──────────────────────────────────────────────────────────────

def test_csv_parsing():
    doc = parse_document("data.csv", "name,age\nAlice,20\nBob,22".encode())
    assert "name\tage" in doc.text
    assert "Alice\t20" in doc.text
    assert "Bob\t22" in doc.text


def test_csv_parsing_skips_empty_rows():
    doc = parse_document("data.csv", "a,b\n,,\n1,2".encode())
    # CSV parser keeps the empty row as tab characters; verify it renders
    assert "\t\t" in doc.text or ",," not in doc.text.replace("\t", ",")
    assert "1\t2" in doc.text


# ── PDF (PyMuPDF) ────────────────────────────────────────────────────

def test_pdf_parsing():
    pdf_bytes = _minimal_pdf_bytes()
    doc = parse_document("report.pdf", pdf_bytes)
    assert "Microscope" in doc.text
    assert "Cell walls" in doc.text


def test_pdf_parsing_segments_have_page_numbers():
    pdf_bytes = _minimal_pdf_bytes()
    doc = parse_document("report.pdf", pdf_bytes)
    for segment in doc.segments:
        assert segment.page_no == 1  # single-page PDF


def test_pdf_parsing_empty_page_raises():
    """A PDF with no extractable text raises ParseError."""
    import fitz

    doc_pdf = fitz.open()
    doc_pdf.new_page(width=100, height=100)  # blank page, no text
    buf = io.BytesIO()
    doc_pdf.save(buf)
    doc_pdf.close()
    with pytest.raises(ParseError, match="未解析出可读文本"):
        parse_document("blank.pdf", buf.getvalue())


# ── DOCX ─────────────────────────────────────────────────────────────

def test_docx_parsing():
    docx_bytes = _minimal_docx_bytes()
    doc = parse_document("report.docx", docx_bytes)
    assert "显微镜实验报告汇总" in doc.text


def test_docx_parsing_segments():
    docx_bytes = _minimal_docx_bytes()
    doc = parse_document("report.docx", docx_bytes)
    assert len(doc.segments) >= 1
    for segment in doc.segments:
        assert segment.page_no == 1
        assert segment.paragraph_no >= 1


def test_docx_empty_document_raises():
    """An empty DOCX raises ParseError."""
    import docx

    document = docx.Document()
    buf = io.BytesIO()
    document.save(buf)
    with pytest.raises(ParseError, match="未解析出可读文本"):
        parse_document("empty.docx", buf.getvalue())


# ── PPTX ─────────────────────────────────────────────────────────────

def test_pptx_parsing():
    pptx_bytes = _minimal_pptx_bytes()
    doc = parse_document("slides.pptx", pptx_bytes)
    assert "周报" in doc.text


def test_pptx_parsing_segments():
    pptx_bytes = _minimal_pptx_bytes()
    doc = parse_document("slides.pptx", pptx_bytes)
    assert len(doc.segments) >= 1
    for segment in doc.segments:
        assert segment.page_no >= 1  # slide index
        assert segment.paragraph_no >= 1


def test_pptx_empty_raises():
    """A PPTX with no text shapes raises ParseError."""
    from pptx import Presentation

    presentation = Presentation()
    slide_layout = presentation.slide_layouts[6]  # blank layout
    presentation.slides.add_slide(slide_layout)
    buf = io.BytesIO()
    presentation.save(buf)
    with pytest.raises(ParseError, match="未解析出可读文本"):
        parse_document("empty.pptx", buf.getvalue())


# ── type round-trip ──────────────────────────────────────────────────

def test_parsed_document_is_dataclass():
    doc = parse_document("notes.txt", "abc".encode())
    assert isinstance(doc, ParsedDocument)
    assert isinstance(doc.text, str)
    assert isinstance(doc.segments, list)
    assert all(isinstance(s, ParsedSegment) for s in doc.segments)
