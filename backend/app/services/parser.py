"""Multi-format document text extraction for the knowledge-base pipeline.

The MVP needs real, searchable text from uploaded PDF/DOCX/PPTX/TXT/MD/CSV
files. Previously the service only ran a raw UTF-8 decode, which yields
garbage for binary formats. This module backs the indexing path with
dependency-driven extractors while keeping them import-light and resilient.

Heavy libraries (PyMuPDF, python-docx, python-pptx, pandas) are imported
lazily inside each extractor so unit tests that exercise only one format stay
fast and do not hard-fail when an optional dependency is absent.
"""
from __future__ import annotations

import io
import re
from dataclasses import dataclass
from pathlib import PurePath


class ParseError(Exception):
    """Raised when a file has no supported extractor or yields no text."""


@dataclass(frozen=True)
class ParsedSegment:
    content: str
    page_no: int = 1
    paragraph_no: int = 1


@dataclass(frozen=True)
class ParsedDocument:
    text: str
    segments: list[ParsedSegment]


def chunk_text(text: str, chunk_size: int = 512, overlap: int = 128) -> list[str]:
    """Split text into overlapping chunks using sliding window (FR-K04)."""
    if not text.strip():
        return []
    chunks: list[str] = []
    start = 0
    text_len = len(text)
    while start < text_len:
        end = min(start + chunk_size, text_len)
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        if end >= text_len:
            break
        start = end - overlap
    return chunks


_FORMAT_ALIASES = {
    "pdf": "pdf", "docx": "docx", "pptx": "pptx",
    "txt": "text", "text": "text",
    "md": "markdown", "markdown": "markdown",
    "csv": "csv",
    "jpg": "image", "jpeg": "image", "png": "image",
    "gif": "image", "bmp": "image", "webp": "image", "tiff": "image",
}


def _format_for(filename: str, declared_type: str | None = None) -> str:
    suffix = PurePath(filename).suffix.lower().lstrip(".")
    resolved = _FORMAT_ALIASES.get(suffix) or _FORMAT_ALIASES.get(declared_type or "")
    if not resolved:
        raise ParseError(f"不支持的解析格式: {suffix or declared_type or 'unknown'}")
    return resolved


def _ocr_fallback(content: bytes, *, is_image: bool = False) -> str:
    """Try OCR on images or scanned PDFs."""
    try:
        from paddleocr import PaddleOCR
        ocr = PaddleOCR(lang="ch", use_angle_cls=False, show_log=False)
        if is_image:
            result = ocr.ocr(content, cls=False)
            if result and result[0]:
                texts = [line[1][0] for line in result[0]]
                return "\n".join(texts) if texts else ""
        else:
            import fitz
            doc = fitz.open(stream=content, filetype="pdf")
            texts = []
            for page in doc:
                pix = page.get_pixmap(dpi=200)
                result = ocr.ocr(pix.tobytes("png"), cls=False)
                if result and result[0]:
                    texts.extend(line[1][0] for line in result[0])
            doc.close()
            if texts: return "\n".join(texts)
    except ImportError: pass

    try:
        import pytesseract
        from PIL import Image
        import io
        img = Image.open(io.BytesIO(content))
        text = pytesseract.image_to_string(img, lang="chi_sim+eng")
        if text.strip(): return text.strip()
    except Exception: pass

    raise ParseError("OCR 未识别出文字（需要安装 paddleocr 或 tesseract）")


def _block_to_segments(text: str, *, base_page: int = 1) -> list[ParsedSegment]:
    raw_lines = text.replace("\r\n", "\n").replace("\r", "\n").split("\n")
    segments: list[ParsedSegment] = []
    paragraph_no = 0
    for line in raw_lines:
        line = line.strip()
        if not line:
            continue
        paragraph_no += 1
        segments.append(ParsedSegment(content=line, page_no=base_page, paragraph_no=paragraph_no))
    if not segments:
        segments.append(ParsedSegment(content=text.strip(), page_no=base_page, paragraph_no=1))
    return segments


def _parse_text(content: bytes) -> str:
    return content.decode("utf-8", errors="replace").strip()


def _parse_markdown(content: bytes) -> str:
    # Keep headings/prose readable; drop fenced code blocks that add noise to chunks.
    text = content.decode("utf-8", errors="replace")
    text = re.sub(r"```.*?```", "", text, flags=re.DOTALL)
    return text.strip()


def _parse_csv(content: bytes) -> str:
    import csv

    text = content.decode("utf-8", errors="replace")
    reader = csv.reader(io.StringIO(text))
    rows = ["\t".join(cell.strip() for cell in row) for row in reader]
    return "\n".join(row for row in rows if row).strip()


def _parse_pdf(content: bytes) -> list[ParsedSegment]:
    import fitz

    doc = fitz.open(stream=content, filetype="pdf")
    try:
        segments: list[ParsedSegment] = []
        paragraph_no = 0
        for page_index, page in enumerate(doc, start=1):
            page_text = page.get_text("text").strip()
            if not page_text:
                continue
            for line in page_text.split("\n"):
                line = line.strip()
                if not line:
                    continue
                paragraph_no += 1
                segments.append(ParsedSegment(content=line, page_no=page_index, paragraph_no=paragraph_no))
    finally:
        doc.close()
    if not segments:
        raise ParseError("PDF 未解析出可读文本（可能是扫描件，需先做 OCR）")
    return segments


def _parse_docx(content: bytes) -> list[ParsedSegment]:
    import docx

    document = docx.Document(io.BytesIO(content))
    segments: list[ParsedSegment] = []
    paragraph_no = 0
    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        paragraph_no += 1
        segments.append(ParsedSegment(content=text, page_no=1, paragraph_no=paragraph_no))
    if not segments:
        raise ParseError("DOCX 未解析出可读文本")
    return segments


def _parse_pptx(content: bytes) -> list[ParsedSegment]:
    from pptx import Presentation

    presentation = Presentation(io.BytesIO(content))
    segments: list[ParsedSegment] = []
    paragraph_no = 0
    for slide_index, slide in enumerate(presentation.slides, start=1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if not text:
                    continue
                paragraph_no += 1
                segments.append(ParsedSegment(content=text, page_no=slide_index, paragraph_no=paragraph_no))
    if not segments:
        raise ParseError("PPTX 未解析出可读文本")
    return segments


def parse_document(
    filename: str,
    content: bytes,
    declared_type: str | None = None,
) -> ParsedDocument:
    """Extract text + segments from a document, honoring its real format."""
    fmt = _format_for(filename, declared_type)
    try:
        if fmt in {"text", "markdown", "csv"}:
            if fmt == "text":
                text = _parse_text(content)
            elif fmt == "markdown":
                text = _parse_markdown(content)
            else:
                text = _parse_csv(content)
            segments = _block_to_segments(text)
        elif fmt == "pdf":
            try:
                segments = _parse_pdf(content)
                text = "\n".join(segment.content for segment in segments)
            except ParseError:
                # FR-K03: try OCR fallback for scanned PDFs
                text = _ocr_fallback(content)
                segments = _block_to_segments(text, base_page=1)
        elif fmt == "docx":
            segments = _parse_docx(content)
            text = "\n".join(segment.content for segment in segments)
        elif fmt == "pptx":
            segments = _parse_pptx(content)
            text = "\n".join(segment.content for segment in segments)
        else:  # pragma: no cover - guarded by _format_for
            raise ParseError(f"未实现的解析格式: {fmt}")
    except ParseError:
        raise
    except Exception as exc:  # noqa: BLE001 - normalize any extractor failure
        raise ParseError(f"解析失败: {exc}") from exc

    if not text.strip():
        raise ParseError("解析结果为空，无可索引文本")
    return ParsedDocument(text=text, segments=segments)
